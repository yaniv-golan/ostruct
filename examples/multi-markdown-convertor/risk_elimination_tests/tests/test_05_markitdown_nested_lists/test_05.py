#!/usr/bin/env python3
"""
Test 5: MarkItDown preserves nested lists depth ≥ 3
PPTX with deeply nested bullet slide
"""

import json
import sys
import subprocess
import tempfile
import os
from pathlib import Path
from typing import Dict, Any, Optional, List
import re


def create_test_pptx() -> Path:
    """Create a test PPTX file with deeply nested bullet points."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        # Create presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = "Deeply Nested List Test"

        # Get content placeholder
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Level 1 Item A"

        # Add nested items
        p = tf.add_paragraph()
        p.text = "Level 2 Item A.1"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Level 3 Item A.1.a"
        p.level = 2

        p = tf.add_paragraph()
        p.text = "Level 4 Item A.1.a.i"
        p.level = 3

        p = tf.add_paragraph()
        p.text = "Level 5 Item A.1.a.i.α"
        p.level = 4

        p = tf.add_paragraph()
        p.text = "Level 3 Item A.1.b"
        p.level = 2

        p = tf.add_paragraph()
        p.text = "Level 2 Item A.2"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Level 1 Item B"
        p.level = 0

        # Save to temp file
        temp_file = Path(tempfile.mktemp(suffix=".pptx"))
        prs.save(str(temp_file))
        return temp_file

    except ImportError:
        # Fallback: use existing PPTX file if available
        test_inputs_dir = Path(__file__).parent.parent.parent / "test-inputs"
        pptx_files = list(test_inputs_dir.glob("*.pptx"))
        if pptx_files:
            return pptx_files[0]
        else:
            raise RuntimeError(
                "python-pptx not available and no test PPTX files found"
            )


def convert_with_markitdown(pptx_path: Path) -> Optional[str]:
    """Convert PPTX to Markdown using MarkItDown."""
    try:
        # Try to import and use markitdown
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert(str(pptx_path))
        return result.text_content

    except ImportError:
        # Fallback to subprocess if markitdown is installed as CLI
        try:
            result = subprocess.run(
                ["markitdown", str(pptx_path)],
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode == 0:
                return result.stdout
            else:
                print(f"MarkItDown CLI error: {result.stderr}")
                return None

        except subprocess.TimeoutExpired:
            print("MarkItDown conversion timed out")
            return None
        except FileNotFoundError:
            print("MarkItDown not found - please install markitdown")
            return None
        except Exception as e:
            print(f"MarkItDown conversion failed: {e}")
            return None


def analyze_nested_lists(markdown_text: str) -> Dict[str, Any]:
    """Analyze nested list depth in markdown text."""
    depth_distribution: Dict[int, int] = {}
    max_depth = 0

    analysis = {
        "has_lists": False,
        "max_depth": max_depth,
        "depth_distribution": depth_distribution,
        "preserves_deep_nesting": False,
    }

    lines = markdown_text.split("\n")

    for line in lines:
        # Count leading spaces/tabs to determine nesting level
        stripped = line.lstrip()
        if stripped.startswith(("- ", "* ", "+ ")) or re.match(
            r"^\d+\.\s", stripped
        ):
            analysis["has_lists"] = True

            # Calculate depth based on indentation
            indent = len(line) - len(stripped)
            # Assume 2 or 4 spaces per level
            depth = (indent // 2) + 1 if indent % 4 != 0 else (indent // 4) + 1

            max_depth = max(max_depth, depth)

            if depth not in depth_distribution:
                depth_distribution[depth] = 0
            depth_distribution[depth] += 1

    # Update analysis with final values
    analysis["max_depth"] = max_depth
    analysis["preserves_deep_nesting"] = max_depth >= 3

    return analysis


def run_test() -> Dict[str, Any]:
    """
    Run test 5: Check if MarkItDown preserves nested lists depth ≥ 3.

    Returns:
        Dict with test results
    """
    results: Dict[str, Any] = {
        "test_id": "05",
        "test_name": "MarkItDown preserves nested lists depth ≥ 3",
        "test_case": "PPTX with deeply nested bullet slide",
        "success": False,
        "error": None,
        "details": {},
    }

    try:
        print(f"Test 5: MarkItDown preserves nested lists depth ≥ 3")
        print(f"Test case: PPTX with deeply nested bullet slide")

        # Create or find test PPTX
        print("Creating test PPTX with deeply nested lists...")
        pptx_path = create_test_pptx()
        results["details"]["pptx_path"] = str(pptx_path)

        # Convert with MarkItDown
        print("Converting PPTX to Markdown with MarkItDown...")
        markdown_text = convert_with_markitdown(pptx_path)

        if markdown_text is None:
            results["error"] = "MarkItDown conversion failed"
            return results

        results["details"]["markdown_output"] = (
            markdown_text[:500] + "..."
            if len(markdown_text) > 500
            else markdown_text
        )

        # Analyze nested lists
        print("Analyzing nested list depth in output...")
        analysis = analyze_nested_lists(markdown_text)
        results["details"]["list_analysis"] = analysis

        # Determine success
        if analysis["has_lists"]:
            if analysis["preserves_deep_nesting"]:
                results["success"] = True
                results["details"]["result"] = (
                    f"PASS: MarkItDown preserves depth {analysis['max_depth']} (≥ 3)"
                )
                print(
                    f"✅ PASS: MarkItDown preserves depth {analysis['max_depth']} (≥ 3)"
                )
            else:
                results["success"] = False
                results["details"]["result"] = (
                    f"FAIL: Max depth only {analysis['max_depth']} (< 3)"
                )
                print(f"❌ FAIL: Max depth only {analysis['max_depth']} (< 3)")
        else:
            results["success"] = False
            results["details"]["result"] = "FAIL: No lists found in output"
            print("❌ FAIL: No lists found in output")

        # Clean up temp file
        if pptx_path.name.startswith("tmp"):
            try:
                pptx_path.unlink()
            except:
                pass

    except Exception as e:
        results["error"] = str(e)
        results["details"]["exception"] = str(e)
        print(f"❌ Test failed with error: {e}")

    return results


def main():
    """Run the test."""
    results = run_test()

    # Save results
    output_file = Path(__file__).parent / "results.json"
    with open(output_file, "w") as f:
        json.dump(results, f, indent=2)

    print(f"Results saved to: {output_file}")
    return results


if __name__ == "__main__":
    main()
