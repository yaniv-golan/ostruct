#!/usr/bin/env python3
"""
Test 8: `python-pptx` retains text order matching visual order
Slide with two textboxes: left/right
"""

import json
import sys
import tempfile
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple


def create_test_pptx() -> Path:
    """Create a test PPTX file with textboxes in specific visual order."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Create presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add left textbox
        left_textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(2), width=Inches(3), height=Inches(1)
        )
        left_textbox.text_frame.text = (
            "LEFT TEXTBOX - This should appear first"
        )

        # Add right textbox
        right_textbox = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(2), width=Inches(3), height=Inches(1)
        )
        right_textbox.text_frame.text = (
            "RIGHT TEXTBOX - This should appear second"
        )

        # Add title for context
        title_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(0.5), width=Inches(7), height=Inches(1)
        )
        title_box.text_frame.text = "Text Order Test: Left then Right"

        # Save to temp file
        temp_file = Path(tempfile.mktemp(suffix=".pptx"))
        prs.save(str(temp_file))
        return temp_file

    except ImportError:
        # Fallback: use existing PPTX file if available
        test_inputs_dir = Path(__file__).parent.parent.parent / "test-inputs"
        pptx_files = list(test_inputs_dir.glob("*.pptx"))
        if pptx_files:
            return pptx_files[0]
        else:
            raise RuntimeError(
                "python-pptx not available and no test PPTX files found"
            )


def analyze_text_order(pptx_path: Path) -> Dict[str, Any]:
    """Analyze text extraction order from PPTX slides."""
    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))

        analysis = {
            "slides_found": len(prs.slides),
            "slide_details": [],
            "text_order_preserved": False,
        }

        for slide_idx, slide in enumerate(prs.slides):
            slide_info = {
                "slide_index": slide_idx,
                "shapes_found": len(slide.shapes),
                "text_shapes": [],
                "extraction_order": [],
            }

            # Extract text from shapes in order
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_content = shape.text_frame.text.strip()
                    if text_content:
                        shape_info = {
                            "shape_index": shape_idx,
                            "text": text_content,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        }
                        slide_info["text_shapes"].append(shape_info)
                        slide_info["extraction_order"].append(text_content)

            # Check if left-to-right order is preserved
            if len(slide_info["text_shapes"]) >= 2:
                # Sort by left position to get expected visual order
                visual_order = sorted(
                    slide_info["text_shapes"], key=lambda x: x["left"]
                )
                visual_texts = [shape["text"] for shape in visual_order]

                # Compare with extraction order
                slide_info["visual_order"] = visual_texts
                slide_info["matches_visual"] = (
                    visual_texts == slide_info["extraction_order"]
                )

                if slide_info["matches_visual"]:
                    analysis["text_order_preserved"] = True

            analysis["slide_details"].append(slide_info)

        return analysis

    except ImportError:
        return {
            "error": "python-pptx not available",
            "slides_found": 0,
            "slide_details": [],
            "text_order_preserved": False,
        }
    except Exception as e:
        return {
            "error": str(e),
            "slides_found": 0,
            "slide_details": [],
            "text_order_preserved": False,
        }


def run_test() -> Dict[str, Any]:
    """
    Run test 8: Check if python-pptx retains text order matching visual order.

    Returns:
        Dict with test results
    """
    results: Dict[str, Any] = {
        "test_id": "08",
        "test_name": "`python-pptx` retains text order matching visual order",
        "test_case": "Slide with two textboxes: left/right",
        "success": False,
        "error": None,
        "details": {},
    }

    try:
        print(
            f"Test 8: `python-pptx` retains text order matching visual order"
        )
        print(f"Test case: Slide with two textboxes: left/right")

        # Create or find test PPTX
        print("Creating test PPTX with left/right textboxes...")
        pptx_path = create_test_pptx()
        results["details"]["pptx_path"] = str(pptx_path)

        # Analyze text order
        print("Analyzing text extraction order...")
        analysis = analyze_text_order(pptx_path)
        results["details"]["text_analysis"] = analysis

        # Determine success
        if "error" in analysis:
            results["error"] = analysis["error"]
            results["success"] = False
            results["details"]["result"] = f"FAIL: {analysis['error']}"
            print(f"❌ FAIL: {analysis['error']}")
        else:
            if analysis["text_order_preserved"]:
                results["success"] = True
                results["details"]["result"] = (
                    "PASS: Text order matches visual left-to-right order"
                )
                print("✅ PASS: Text order matches visual left-to-right order")
            else:
                results["success"] = False
                results["details"]["result"] = (
                    "FAIL: Text order does not match visual order"
                )
                print("❌ FAIL: Text order does not match visual order")

                # Show details
                for slide in analysis["slide_details"]:
                    if "visual_order" in slide and "extraction_order" in slide:
                        print(f"   Visual order: {slide['visual_order']}")
                        print(
                            f"   Extraction order: {slide['extraction_order']}"
                        )

        # Clean up temp file
        if pptx_path.name.startswith("tmp"):
            try:
                pptx_path.unlink()
            except:
                pass

    except Exception as e:
        results["error"] = str(e)
        results["details"]["exception"] = str(e)
        print(f"❌ Test failed with error: {e}")

    return results


def main():
    """Run the test."""
    results = run_test()

    # Save results
    output_file = Path(__file__).parent / "results.json"
    with open(output_file, "w") as f:
        json.dump(results, f, indent=2)

    print(f"Results saved to: {output_file}")
    return results


if __name__ == "__main__":
    main()
